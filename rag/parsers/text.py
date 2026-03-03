"""Plain text, Markdown, HTML, CSV, JSON, YAML parser."""

from __future__ import annotations

import json
import re
from pathlib import Path


def parse(path: str) -> str:
    """
    Parse a text-based file to plain text.

    Handles: .txt, .md, .html, .htm, .csv, .json, .yaml, .yml, .pptx

    Args:
        path: File path string.

    Returns:
        Extracted plain text content.
    """
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix in (".html", ".htm"):
        return _parse_html(p)
    elif suffix == ".json":
        return _parse_json(p)
    elif suffix in (".yaml", ".yml"):
        return _parse_yaml(p)
    elif suffix == ".pptx":
        return _parse_pptx(p)
    else:
        return p.read_text(encoding="utf-8", errors="replace")


def _parse_html(p: Path) -> str:
    """Strip HTML tags from an HTML file."""
    html = p.read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"<(script|style)[^>]*>.*?</\1>", "", html, flags=re.DOTALL | re.I)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    return text


def _parse_json(p: Path) -> str:
    """Convert JSON to readable text."""
    try:
        data = json.loads(p.read_text(encoding="utf-8"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return p.read_text(encoding="utf-8", errors="replace")


def _parse_yaml(p: Path) -> str:
    """Convert YAML to readable text."""
    try:
        import yaml

        data = yaml.safe_load(p.read_text(encoding="utf-8"))
        return json.dumps(data, indent=2, ensure_ascii=False)
    except Exception:
        return p.read_text(encoding="utf-8", errors="replace")


def _parse_pptx(p: Path) -> str:
    """Extract text from a PowerPoint file."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped]

        prs = Presentation(str(p))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n\n".join(texts)
    except ImportError:
        return f"[python-pptx not installed, cannot parse {p.name}]"
    except Exception as exc:
        return f"[PPTX parse error: {exc}]"
